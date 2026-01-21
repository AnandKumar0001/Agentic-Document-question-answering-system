"""Document loader for handling various file formats."""
import os
from pathlib import Path
from typing import List, Dict, Any
from pypdf import PdfReader
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
import pandas as pd
from pptx import Presentation
from datetime import datetime


class DocumentLoader:
    """Loads and extracts content from various document types."""

    def __init__(self):
        self.supported_formats = {".pdf", ".txt", ".csv", ".xlsx", ".png", ".jpg", ".jpeg", ".pptx"}

    def load_documents(self, file_path: str) -> List[Dict[str, Any]]:
        """Load documents from a file."""
        file_ext = Path(file_path).suffix.lower()

        if file_ext == ".pdf":
            return self._load_pdf(file_path)
        elif file_ext == ".txt":
            return self._load_text(file_path)
        elif file_ext in {".csv", ".xlsx"}:
            return self._load_tabular(file_path)
        elif file_ext in {".png", ".jpg", ".jpeg"}:
            return self._load_image(file_path)
        elif file_ext == ".pptx":
            return self._load_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")

    def _load_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Load and extract content from PDF files."""
        documents = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                documents.append({
                    "content": text,
                    "source": file_path,
                    "page": page_num + 1,
                    "type": "text",
                    "timestamp": datetime.now().isoformat()
                })

        # Extract images from PDF
        try:
            images = convert_from_path(file_path)
            for img_num, image in enumerate(images):
                extracted_text = pytesseract.image_to_string(image)
                documents.append({
                    "content": extracted_text,
                    "source": file_path,
                    "page": img_num + 1,
                    "type": "image",
                    "timestamp": datetime.now().isoformat()
                })
        except Exception as e:
            print(f"Warning: Could not extract images from PDF: {e}")

        return documents

    def _load_text(self, file_path: str) -> List[Dict[str, Any]]:
        """Load text files."""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        return [{
            "content": content,
            "source": file_path,
            "type": "text",
            "timestamp": datetime.now().isoformat()
        }]

    def _load_tabular(self, file_path: str) -> List[Dict[str, Any]]:
        """Load CSV and Excel files."""
        if file_path.endswith('.csv'):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)

        # Convert dataframe to markdown table format
        table_text = df.to_markdown(index=False)
        
        return [{
            "content": table_text,
            "source": file_path,
            "type": "table",
            "timestamp": datetime.now().isoformat(),
            "rows": len(df),
            "columns": list(df.columns)
        }]

    def _load_image(self, file_path: str) -> List[Dict[str, Any]]:
        """Load and extract text from images using OCR."""
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)

        return [{
            "content": text,
            "source": file_path,
            "type": "image",
            "timestamp": datetime.now().isoformat()
        }]

    def _load_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """Load PowerPoint presentations."""
        presentation = Presentation(file_path)
        documents = []

        for slide_num, slide in enumerate(presentation.slides):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            documents.append({
                "content": "\n".join(slide_text),
                "source": file_path,
                "slide": slide_num + 1,
                "type": "presentation",
                "timestamp": datetime.now().isoformat()
            })

        return documents

    def load_batch(self, directory: str) -> List[Dict[str, Any]]:
        """Load all supported documents from a directory."""
        all_documents = []
        directory_path = Path(directory)

        for file_path in directory_path.rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in self.supported_formats:
                try:
                    docs = self.load_documents(str(file_path))
                    all_documents.extend(docs)
                except Exception as e:
                    print(f"Error loading {file_path}: {e}")

        return all_documents
